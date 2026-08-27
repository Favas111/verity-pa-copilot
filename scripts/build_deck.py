#!/usr/bin/env python3
"""
Build the prototype submission deck.

    python3 scripts/build_deck.py [out.pptx]

Uses the official template's own background art (extracted from the supplied
.pptx) so the deck carries hackathon branding, while giving us full control
of the content. Slide size matches the template exactly: 10 x 5.625in
(LAYOUT_16x9), NOT the 13.33in wide layout — coordinates past 10in are
silently written off-slide.

Content maps to the three sections the template requires:
  1. Problem Brief        -> slides 2-3
  2. Architecture Diagram -> slide 4 (names the CoCo CLI skills explicitly)
  3. Impact Statement     -> slide 7

Every figure on the Impact slide is measured from this build, not quoted
from industry sources. See docs/metrics.md.
"""

import os
import subprocess
import sys
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
TEMPLATE = os.path.expanduser(
    "~/Downloads/Prototype Submission Template _ CoCo CLI Hackathon GCC Edition.pptx"
)

# ---------------------------------------------------------------------
# Team details — REPLACE BEFORE SUBMITTING
# ---------------------------------------------------------------------
TEAM_NAME = "Coyot"
TEAM_LEADER = "Priyanka Karmakar"
TEAM_SIZE = "2 — Priyanka Karmakar (lead), Mohammed Favas"
PROBLEM_STATEMENT = "Patient & Member 360 / Clinical-Regulatory Copilot (HCLS)"

# ---------------------------------------------------------------------
# Palette — brand blue from the hackathon logo, plus semantic verdict
# colours kept deliberately distinct from it so "needs attention" never
# reads as decoration.
# ---------------------------------------------------------------------
BRAND = RGBColor(0x1B, 0xA0, 0xF2)
NAVY = RGBColor(0x0B, 0x1A, 0x2B)
INK = RGBColor(0x1A, 0x1A, 0x1A)
BODY = RGBColor(0x33, 0x40, 0x45)
MUTED = RGBColor(0x6B, 0x7D, 0x84)
RULE = RGBColor(0xD8, 0xE0, 0xE3)
ARROW = RGBColor(0x9A, 0xAB, 0xB2)   # flow arrows; RULE is too faint to read
MET = RGBColor(0x1B, 0x6E, 0x45)
MET_BG = RGBColor(0xE2, 0xF0, 0xE8)
GAP = RGBColor(0x8A, 0x5A, 0x06)
GAP_BG = RGBColor(0xF7, 0xED, 0xDA)
STOP = RGBColor(0x9C, 0x31, 0x29)
STOP_BG = RGBColor(0xF8, 0xE7, 0xE5)
WASH = RGBColor(0xF2, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SANS = "Arial"
MONO = "Courier New"

# Safe content band: below the black header bar, above the gradient footer.
TOP = 0.66
BOTTOM = 5.34
LEFT = 0.46
RIGHT = 9.54
W = RIGHT - LEFT


def extract_backgrounds(outdir):
    os.makedirs(outdir, exist_ok=True)
    paths = {}
    with zipfile.ZipFile(TEMPLATE) as z:
        for key, name in (("title", "image5.png"), ("content", "image4.png")):
            p = os.path.join(outdir, name)
            with open(p, "wb") as f:
                f.write(z.read(f"ppt/media/{name}"))
            paths[key] = p
    return paths


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, size=10, bold=False, color=BODY, font=SANS, space_after=4,
         first=False, align=None, space_before=0, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    # Autoshape text frames default to CENTER in PowerPoint. Left is the
    # intent everywhere except chips, and mixing the two within one card
    # (centered heading over left body) is the tell of a careless deck.
    p.alignment = PP_ALIGN.LEFT if align is None else align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def rich(tf, parts, size=10, space_after=4, first=False, align=None,
         space_before=0):
    """parts = [(text, {bold, color, font, italic, size}), ...]"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT if align is None else align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    for text, style in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(style.get("size", size))
        r.font.bold = style.get("bold", False)
        r.font.italic = style.get("italic", False)
        r.font.color.rgb = style.get("color", BODY)
        r.font.name = style.get("font", SANS)
    return p


def box(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        line_w=0.75, adj=0.10):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = adj
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    return s, tf


def add_slide(prs, bg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_picture(bg, 0, 0, width=prs.slide_width,
                             height=prs.slide_height)
    return slide


def title_block(slide, title, kicker=None, sub=None):
    y = TOP + 0.10
    if kicker:
        tf = textbox(slide, LEFT, y, W, 0.20)
        para(tf, kicker.upper(), size=8, bold=True, color=BRAND, first=True,
             space_after=0)
        y += 0.22
    tf = textbox(slide, LEFT, y, W, 0.42)
    para(tf, title, size=21, bold=True, color=NAVY, first=True, space_after=0)
    y += 0.40
    if sub:
        tf = textbox(slide, LEFT, y, W, 0.30)
        para(tf, sub, size=10, color=MUTED, first=True, space_after=0)
        y += 0.32
    return y + 0.14


def chip(slide, x, y, text, fg, bg, w=0.86, h=0.20, size=7.5):
    s, tf = box(slide, x, y, w, h, fill=bg, line=None, adj=0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, text, size=size, bold=True, color=fg, font=MONO, first=True,
         space_after=0, align=PP_ALIGN.CENTER)
    return s


# =====================================================================
# Slides
# =====================================================================

def slide_title(prs, bg):
    slide = add_slide(prs, bg["title"])
    # White area begins around 52% down; the hero art occupies the top.
    y = 3.34
    tf = textbox(slide, LEFT, y, W, 0.34)
    rich(tf, [("VERITY", {"bold": True, "color": NAVY, "size": 25}),
              ("  ", {}),
              ("Prior Authorization Evidence Copilot",
               {"color": MUTED, "size": 14})],
         first=True, space_after=0)

    tf = textbox(slide, LEFT, y + 0.44, W, 0.24)
    para(tf, "Every determination traces to a cited source. The system can approve, "
             "or route to a clinician — it cannot deny.",
         size=9.5, color=BODY, first=True, space_after=0, italic=True)

    rows = [("Team Name", TEAM_NAME),
            ("Problem Statement", PROBLEM_STATEMENT),
            ("Team Leader Name", TEAM_LEADER),
            ("Team Size", TEAM_SIZE)]
    yy = y + 0.86
    for label, value in rows:
        tf = textbox(slide, LEFT, yy, 1.55, 0.20)
        para(tf, label, size=9, bold=True, color=NAVY, first=True, space_after=0)
        tf = textbox(slide, LEFT + 1.62, yy, W - 1.62, 0.20)
        para(tf, value, size=9, color=BODY, first=True, space_after=0)
        yy += 0.27


def slide_problem(prs, bg):
    slide = add_slide(prs, bg["content"])
    y = title_block(
        slide, "A reviewer cannot see the evidence that would approve the case",
        kicker="Problem Brief",
        sub="Prior authorization for specialty drugs — Meridian Health Plan (synthetic)")

    # Persona + pain, two columns
    cw = (W - 0.22) / 2
    s, tf = box(slide, LEFT, y, cw, 1.62, fill=WHITE, line=RULE)
    para(tf, "THE USER", size=7.5, bold=True, color=BRAND, first=True, space_after=5)
    para(tf, "Utilization-management nurse", size=11, bold=True, color=NAVY,
         space_after=4)
    para(tf, "Decides whether a requested drug meets the plan's published medical "
             "policy, then documents why. Works a queue under a turnaround clock.",
         size=9, color=BODY, space_after=0)

    s, tf = box(slide, LEFT + cw + 0.22, y, cw, 1.62, fill=WHITE, line=RULE)
    para(tf, "THE PAIN", size=7.5, bold=True, color=BRAND, first=True, space_after=5)
    para(tf, "The evidence is split in half", size=11, bold=True, color=NAVY,
         space_after=4)
    para(tf, "Eligibility, claims, labs and pharmacy fills live in tables. Trial "
             "history, intolerances and contraindications live in narrative notes "
             "and scanned faxes. The policy itself is a PDF. Nothing joins them but "
             "a person reading.",
         size=9, color=BODY, space_after=0)

    y2 = y + 1.78
    s, tf = box(slide, LEFT, y2, W, 0.92, fill=STOP_BG, line=None)
    para(tf, "WHAT GOES WRONG", size=7.5, bold=True, color=STOP, first=True,
         space_after=5)
    rich(tf, [("A request is refused because the evidence was never found — not "
               "because it did not exist. ", {"color": INK, "size": 9.5}),
              ("A metformin trial completed under a previous carrier, or an "
               "intolerance recorded by an out-of-network specialist, is invisible "
               "to a query scoped to the current plan.",
               {"color": INK, "size": 9.5})],
         space_after=0)

    y3 = y2 + 1.06
    # Figures verified against CMS's own fact sheet for CMS-0057-F: 72 hours
    # expedited, 7 calendar days standard, a specific denial reason required,
    # in force for impacted payers from 1 January 2026.
    tf = textbox(slide, LEFT, y3, W, 0.50)
    rich(tf, [("Domain context.  ", {"bold": True, "color": NAVY, "size": 9}),
              ("Since January 2026, CMS-0057-F requires impacted payers to return "
               "prior-authorization decisions within 72 hours (expedited) or 7 calendar "
               "days (standard), and to give a specific reason for every denial. "
               "\"The model decided\" is not a reason. Each criterion must be "
               "individually defensible and traceable to a source.",
               {"color": BODY, "size": 9})],
         first=True, space_after=0)


def slide_concept(prs, bg):
    slide = add_slide(prs, bg["content"])
    y = title_block(
        slide, "The model reads evidence. Arithmetic renders the verdict.",
        kicker="The Criteria Ledger",
        sub="A policy PDF becomes a tree of individually checkable nodes")

    # Tree illustration, left
    tw = 4.62
    s, tf = box(slide, LEFT, y, tw, 2.62, fill=WHITE, line=RULE)
    para(tf, "MHP-PA-0142  ·  21 NODES", size=7, bold=True, color=BRAND,
         font=MONO, first=True, space_after=6)

    tree = [
        (0, "[ALL_OF] Medical necessity", None),
        (1, "[ALL_OF] §3 Coverage criteria", None),
        (2, "§3.1 Age ≥ 18", "structured"),
        (2, "§3.3 HbA1c ≥ 7.0% within 90d", "structured"),
        (1, "[ALL_OF] §4 Step therapy", None),
        (2, "[ANY_OF] §4.1 Metformin", None),
        (3, "adequate trial", "structured"),
        (3, "documented intolerance", "narrative"),
        (1, "[NONE_OF] §5 Exclusions", None),
        (2, "§5.1 MTC / MEN2 history", "narrative"),
    ]
    for depth, text, kind in tree:
        pad = "   " * depth
        col = NAVY if kind is None else BODY
        parts = [(pad + text, {"font": MONO, "size": 8,
                               "bold": kind is None, "color": col})]
        if kind:
            parts.append(("  (" + kind + ")",
                          {"font": SANS, "size": 7, "color": MUTED}))
        rich(tf, parts, space_after=2)

    # Right column: how it resolves
    rx = LEFT + tw + 0.22
    rw = W - tw - 0.22
    steps = [
        ("Each leaf resolves alone",
         "A structured leaf is SQL over claims, labs and fills. A narrative leaf "
         "is retrieval over clinical notes. Each returns MET, NOT MET, or NO "
         "EVIDENCE — and its citation."),
        ("Groups fold upward",
         "ALL_OF, ANY_OF and NONE_OF combine children in deterministic SQL. "
         "No model call decides an outcome."),
        ("Two outcomes exist",
         "APPROVE, or ROUTE TO CLINICIAN with the missing item named. There is no "
         "code path that denies — a structural property, not a prompt instruction."),
    ]
    yy = y
    for i, (head, text) in enumerate(steps):
        h = 0.86
        s, tf = box(slide, rx, yy, rw, h, fill=WASH, line=None)
        para(tf, head, size=9.5, bold=True, color=NAVY, first=True, space_after=3)
        para(tf, text, size=8.5, color=BODY, space_after=0)
        yy += h + 0.02

    tf = textbox(slide, LEFT, y + 2.72, W, 0.28)
    rich(tf, [("Why it matters.  ", {"bold": True, "color": NAVY, "size": 8.5}),
              ("A verdict computed from cited leaves cannot be uncited, and cannot "
               "drift between runs on the same evidence.",
               {"color": BODY, "size": 8.5})],
         first=True, space_after=0)


def slide_architecture(prs, bg):
    slide = add_slide(prs, bg["content"])
    y = title_block(slide, "Architecture", kicker="How it fits together",
                    sub="Everything executes inside Snowflake. No PHI leaves the perimeter.")

    lane_h = 0.60
    gap = 0.13
    bw = 1.66
    ax = LEFT

    def flow(yy, label, label_col, items):
        tf = textbox(slide, ax, yy + 0.14, 0.86, 0.3)
        para(tf, label, size=7.5, bold=True, color=label_col, first=True,
             space_after=0)
        x = ax + 0.92
        for i, (title, sub, style) in enumerate(items):
            fill, edge, tcol = style
            s, tf = box(slide, x, yy, bw, lane_h, fill=fill, line=edge)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            para(tf, title, size=8.5, bold=True, color=tcol, first=True,
                 space_after=1)
            para(tf, sub, size=6.5, color=MUTED, font=MONO, space_after=0)
            if i < len(items) - 1:
                ar = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.015),
                    Inches(yy + lane_h / 2 - 0.045), Inches(0.10), Inches(0.09))
                ar.fill.solid()
                ar.fill.fore_color.rgb = ARROW
                ar.line.fill.background()
                ar.shadow.inherit = False
            x += bw + 0.13

    coco = (RGBColor(0xE8, 0xF4, 0xFD), BRAND, NAVY)
    plain = (WHITE, RULE, NAVY)

    flow(y, "POLICY", BRAND, [
        ("Policy PDF", "@POLICY_DOCS", plain),
        ("AI_PARSE_DOCUMENT", "layout mode", plain),
        ("policy-criteria-extractor", "CoCo CLI skill", coco),
        ("Criteria Ledger", "21 nodes", plain),
    ])

    flow(y + lane_h + gap, "MEMBER", BRAND, [
        ("Claims · Labs · Rx", "5,003 members", plain),
        ("V_DRUG_TRIAL", "gap analysis", plain),
        ("Structured leaves", "SQL + citation", plain),
    ])

    flow(y + 2 * (lane_h + gap), "NARRATIVE", BRAND, [
        ("Clinical notes", "CLINICAL_NOTE", plain),
        ("Cortex Search", "member-scoped", plain),
        ("AI_FILTER adjudication", "affirms?", plain),
        ("Narrative leaves", "+ citation", plain),
    ])

    yb = y + 3 * (lane_h + gap) + 0.06
    s, tf = box(slide, LEFT + 0.92, yb, bw * 2 + 0.13, lane_h, fill=NAVY, line=None)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "Deterministic rollup", size=9, bold=True, color=WHITE, first=True,
         space_after=1)
    para(tf, "ALL_OF / ANY_OF / NONE_OF — SQL, not a model call", size=6.5,
         color=RGBColor(0x9F, 0xC4, 0xDE), font=MONO, space_after=0)

    x2 = LEFT + 0.92 + bw * 2 + 0.13 + 0.13
    s, tf = box(slide, x2, yb, bw, lane_h, fill=WHITE, line=RULE)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "AUDIT trail", size=8.5, bold=True, color=NAVY, first=True,
         space_after=1)
    para(tf, "every node + source", size=6.5, color=MUTED, font=MONO, space_after=0)

    s, tf = box(slide, x2 + bw + 0.13, yb, bw, lane_h, fill=WHITE, line=RULE)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "Streamlit console", size=8.5, bold=True, color=NAVY, first=True,
         space_after=1)
    para(tf, "in Snowflake", size=6.5, color=MUTED, font=MONO, space_after=0)

    tf = textbox(slide, LEFT, yb + lane_h + 0.10, W, 0.30)
    rich(tf, [("CoCo CLI skills.  ", {"bold": True, "color": NAVY, "size": 8.5}),
              ("policy-criteria-extractor", {"font": MONO, "color": BRAND, "size": 8.5}),
              (" converts parsed policy prose into the criteria tree, scored blind "
               "against a held-out answer key. It recovered every node independently "
               "and added two the hand-authored key was missing.",
               {"color": BODY, "size": 8.5})],
         first=True, space_after=0)


def slide_hero(prs, bg):
    slide = add_slide(prs, bg["content"])
    y = title_block(
        slide, "The evidence that approves the case is in the wrong place",
        kicker="Demo · member M09000001",
        sub="§4.1 requires three consecutive months of metformin. Hers sits under a prior carrier's member id.")

    cw = (W - 0.22) / 2

    s, tf = box(slide, LEFT, y, cw, 1.30, fill=WHITE, line=RULE)
    para(tf, "CONVENTIONAL QUERY", size=7.5, bold=True, color=MUTED, first=True,
         space_after=4)
    para(tf, "0 fills", size=22, bold=True, color=STOP, space_after=3)
    para(tf, "Scoped to the current member id, step therapy fails. The request is "
             "refused.", size=8.5, color=BODY, space_after=4)
    para(tf, "WHERE member_id = 'M09000001'", size=7, color=MUTED, font=MONO,
         space_after=0)

    s, tf = box(slide, LEFT + cw + 0.22, y, cw, 1.30, fill=WHITE, line=BRAND,
                line_w=1.25)
    para(tf, "VERITY — IDENTITY-LINKED", size=7.5, bold=True, color=BRAND,
         first=True, space_after=4)
    para(tf, "5.9 months", size=22, bold=True, color=MET, space_after=3)
    para(tf, "Six consecutive fills recovered from prior coverage. §4.1 satisfied, "
             "with a citation.", size=8.5, color=BODY, space_after=4)
    para(tf, "source: Northstar Mutual Health", size=7, color=MUTED, font=MONO,
         space_after=0)

    y2 = y + 1.44
    s, tf = box(slide, LEFT, y2, W, 0.56, fill=WASH, line=None)
    rich(tf, [("This is the policy applied correctly, not a loophole.  ",
               {"bold": True, "color": NAVY, "size": 9}),
              ("§2.3 states that a trial completed under a prior member identifier "
               "or prior carrier SHALL count toward step therapy. Conventional "
               "tooling simply never looks there.",
               {"color": BODY, "size": 9})],
         first=True, space_after=0)

    y3 = y2 + 0.70
    tf = textbox(slide, LEFT, y3, W, 0.22)
    para(tf, "THE SAME REVIEW, THREE MEMBERS", size=7.5, bold=True, color=BRAND,
         first=True, space_after=0)

    y4 = y3 + 0.26
    rows = [
        ("Elena Vasquez", "APPROVE", MET, MET_BG,
         "§4.1 rescued by prior-coverage linkage; §4.2 by an out-of-network note "
         "documenting intolerance to both second-line classes."),
        ("Marcus Thorne", "ROUTE", GAP, GAP_BG,
         "Only HbA1c on file is 210 days old. The system names exactly what is "
         "missing rather than refusing."),
        ("Priya Nakamura", "ROUTE", STOP, STOP_BG,
         "Every criterion passes, but a family history of medullary thyroid "
         "carcinoma is buried in a consult note. §5.1 fires."),
    ]
    yy = y4
    for name, verdict, fg, bgc, text in rows:
        tf = textbox(slide, LEFT, yy + 0.02, 1.28, 0.20)
        para(tf, name, size=8.5, bold=True, color=NAVY, first=True, space_after=0)
        chip(slide, LEFT + 1.32, yy, verdict, fg, bgc, w=0.74, h=0.19, size=7)
        tf = textbox(slide, LEFT + 2.16, yy + 0.01, W - 2.16, 0.34)
        para(tf, text, size=8, color=BODY, first=True, space_after=0)
        yy += 0.42


def slide_trust(prs, bg):
    slide = add_slide(prs, bg["content"])
    y = title_block(
        slide, "Retrieval finds the topic. It does not know the answer.",
        kicker="A failure we found and fixed",
        sub="Searching two members for a history of medullary thyroid carcinoma returns a confident hit for both")

    hdr = ["MEMBER", "WHAT RETRIEVAL RETURNED", "NAIVE", "VERITY"]
    widths = [1.30, 4.86, 1.16, 1.22]
    xs = [LEFT]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)

    for i, h in enumerate(hdr):
        tf = textbox(slide, xs[i], y, widths[i], 0.18)
        para(tf, h, size=7, bold=True, color=MUTED, font=MONO, first=True,
             space_after=0)

    yy = y + 0.24
    rows = [
        ("M09000003",
         "\"Mother was diagnosed with medullary thyroid carcinoma at age 52\"",
         "EXCLUDED", STOP, STOP_BG, "EXCLUDED", STOP, STOP_BG, "correct"),
        ("M09000002",
         "\"No personal or family history of thyroid malignancy\"",
         "EXCLUDED", STOP, STOP_BG, "NOT EXCLUDED", MET, MET_BG, "correct"),
    ]
    for mid, text, nv, nfg, nbg, vv, vfg, vbg, note in rows:
        tf = textbox(slide, xs[0], yy + 0.04, widths[0], 0.20)
        para(tf, mid, size=8, color=NAVY, font=MONO, first=True, space_after=0)
        tf = textbox(slide, xs[1], yy + 0.03, widths[1] - 0.12, 0.24)
        para(tf, text, size=8.5, color=BODY, first=True, space_after=0, italic=True)
        chip(slide, xs[2], yy, nv, nfg, nbg, w=0.98, h=0.20, size=6.5)
        chip(slide, xs[3], yy, vv, vfg, vbg, w=1.12, h=0.20, size=6.5)
        yy += 0.40

    s, tf = box(slide, LEFT, yy + 0.06, W, 0.74, fill=STOP_BG, line=None)
    rich(tf, [("Treating a retrieval hit as a met criterion would have blocked care "
               "for a member whose own record rules the exclusion out.",
               {"bold": True, "color": STOP, "size": 9.5})],
         first=True, space_after=3)
    para(tf, "Reranker scores did separate the two, but a score threshold is "
             "arbitrary and drifts with corpus and phrasing. The fix is structural: "
             "Cortex Search narrows to candidates, then a separate adjudication step "
             "decides whether the passage affirmatively asserts the condition. "
             "Retrieval proposes; adjudication disposes.",
         size=8.5, color=INK, space_after=0)

    y3 = yy + 0.90
    cw = (W - 0.20) / 2
    s, tf = box(slide, LEFT, y3, cw, 0.72, fill=WHITE, line=RULE)
    para(tf, "BLIND EXTRACTION, SCORED", size=7, bold=True, color=BRAND,
         font=MONO, first=True, space_after=3)
    rich(tf, [("21/21", {"bold": True, "color": MET, "size": 13}),
              ("  structural match · ", {"color": BODY, "size": 8.5}),
              ("0/21", {"bold": True, "color": NAVY, "size": 13}),
              ("  identical labels", {"color": BODY, "size": 8.5})],
         space_after=0)

    s, tf = box(slide, LEFT + cw + 0.20, y3, cw, 0.72, fill=WHITE, line=RULE)
    para(tf, "THE FIRST RUN WAS DISCARDED", size=7, bold=True, color=BRAND,
         font=MONO, first=True, space_after=3)
    para(tf, "It scored 100% by reading the answer key sitting in the table it was "
             "writing to. We moved the key to a held-out schema and added a "
             "contamination detector.", size=8, color=BODY, space_after=0)


def slide_impact(prs, bg):
    slide = add_slide(prs, bg["content"])
    y = title_block(slide, "Impact", kicker="Impact Statement",
                    sub="Every figure below is measured from this build, not quoted from industry sources")

    # Sub-lines are separate paragraphs: a "\n" inside a python-pptx run does
    # not become a line break, it renders as a stray glyph or collapses.
    metrics = [
        ("43s", "end-to-end adjudication",
         ["21 criteria, 6 retrievals,", "20 adjudication calls"]),
        ("100%", "of determinations cited",
         ["315 audit nodes, each", "with its source"]),
        ("21/21", "blind extraction match",
         ["policy PDF → criteria tree,", "scored against held-out key"]),
        ("$6.77", "total build cost",
         ["245,291 rows,", "2.26 credits"]),
    ]
    mw = (W - 3 * 0.14) / 4
    for i, (big, label, sublines) in enumerate(metrics):
        x = LEFT + i * (mw + 0.14)
        s, tf = box(slide, x, y, mw, 1.02, fill=WHITE, line=RULE)
        para(tf, big, size=20, bold=True, color=BRAND, first=True, space_after=1)
        para(tf, label, size=8.5, bold=True, color=NAVY, space_after=3)
        for ln in sublines:
            para(tf, ln, size=7, color=MUTED, space_after=0)

    y2 = y + 1.14
    tf = textbox(slide, LEFT, y2, W, 0.20)
    para(tf, "WHAT CHANGES FOR THE REVIEWER", size=7.5, bold=True, color=BRAND,
         first=True, space_after=0)

    y3 = y2 + 0.24
    cw = (W - 0.20) / 2
    s, tf = box(slide, LEFT, y3, cw, 0.98, fill=WASH, line=None)
    para(tf, "The evidence packet is assembled, not hunted", size=9.5, bold=True,
         color=NAVY, first=True, space_after=3)
    para(tf, "A reviewer opens a request and sees each criterion with its verdict "
             "and source already attached — including narrative evidence a "
             "structured query would never surface. Clear-cut cases approve "
             "themselves; the rest arrive ready to judge.",
         size=8.5, color=BODY, space_after=0)

    s, tf = box(slide, LEFT + cw + 0.20, y3, cw, 0.98, fill=WASH, line=None)
    para(tf, "The reason is auditable by construction", size=9.5, bold=True,
         color=NAVY, first=True, space_after=3)
    para(tf, "Every determination stores its full node trail, the passage behind "
             "each leaf, and the policy version in force on the date of service. "
             "A CMS-0057-F reason code is a query, not a reconstruction.",
         size=8.5, color=BODY, space_after=0)

    y4 = y3 + 1.12
    tf = textbox(slide, LEFT, y4, W, 0.20)
    para(tf, "BEYOND THE DEMO", size=7.5, bold=True, color=BRAND, first=True,
         space_after=0)

    y5 = y4 + 0.24
    items = [
        ("Domain-agnostic", "The Criteria Ledger is a governed evidence engine. Swap "
                            "medical policy for AML rules or supplier contracts — the "
                            "rollup is unchanged."),
        ("Push upstream", "Share the criteria engine with providers so a request is "
                          "checked before it is submitted, removing the denial "
                          "rather than appealing it."),
        ("Scales by policy", "Adding a policy is adding a PDF. Parsing, extraction "
                             "and scoring already run unattended."),
    ]
    iw = (W - 2 * 0.14) / 3
    for i, (head, text) in enumerate(items):
        x = LEFT + i * (iw + 0.14)
        s, tf = box(slide, x, y5, iw, 0.78, fill=WHITE, line=RULE)
        para(tf, head, size=8.5, bold=True, color=NAVY, first=True, space_after=2)
        para(tf, text, size=7.5, color=BODY, space_after=0)


def main():
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.join(ROOT, "docs", "verity-prototype.pptx")
    bg = extract_backgrounds(os.path.join(ROOT, "docs", "_bg"))

    prs = Presentation()
    prs.slide_width = Emu(9144000)    # 10.0in — matches the template exactly
    prs.slide_height = Emu(5143500)   # 5.625in

    slide_title(prs, bg)
    slide_problem(prs, bg)
    slide_concept(prs, bg)
    slide_architecture(prs, bg)
    slide_hero(prs, bg)
    slide_trust(prs, bg)
    slide_impact(prs, bg)

    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print("wrote", out, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
